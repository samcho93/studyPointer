# -*- coding: utf-8 -*-
"""
C언어 포인터 완전정복 - PowerPoint 프레젠테이션 생성 스크립트
64비트 시스템 기반 메모리 시각화 학습
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# COLOR CONSTANTS
# ============================================================
DARK_NAVY = RGBColor(0x0F, 0x17, 0x2A)
BLUE_PRIMARY = RGBColor(0x1E, 0x40, 0xAF)
BLUE_LIGHT = RGBColor(0x3B, 0x82, 0xF6)
BLUE_ACCENT = RGBColor(0x60, 0xA5, 0xFA)
BLUE_PALE = RGBColor(0xDB, 0xEA, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_BG = RGBColor(0x1E, 0x29, 0x3B)
GRAY_TEXT = RGBColor(0x94, 0xA3, 0xB8)
GRAY_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
GRAY_DARK = RGBColor(0x33, 0x33, 0x33)
GREEN_CODE = RGBColor(0x4A, 0xDE, 0x80)
YELLOW_CODE = RGBColor(0xFD, 0xE0, 0x47)
ORANGE_CODE = RGBColor(0xFB, 0x92, 0x3C)
CYAN_CODE = RGBColor(0x22, 0xD3, 0xEE)

# Memory map row colors
MEM_HEADER_BG = RGBColor(0x1E, 0x3A, 0x5F)
MEM_INT_BG = RGBColor(0xDB, 0xEA, 0xFE)      # blue - int
MEM_PTR_BG = RGBColor(0xD1, 0xFA, 0xE5)       # green - pointer
MEM_CHAR_BG = RGBColor(0xFE, 0xF9, 0xC3)      # yellow - char
MEM_ARR_BG = RGBColor(0xE9, 0xD5, 0xFF)        # purple - array
MEM_HEAP_BG = RGBColor(0xFE, 0xE2, 0xE2)       # red - heap
MEM_FUNC_BG = RGBColor(0xFF, 0xED, 0xD5)       # orange - function/code
MEM_DEFAULT_BG = RGBColor(0xF1, 0xF5, 0xF9)    # light gray - default

FONT_KOREAN = 'Malgun Gothic'
FONT_CODE = 'Consolas'

# ============================================================
# LECTURE DATA - All 20 Lectures
# ============================================================
LECTURES = [
    {
        "id": 1,
        "title": "메모리의 구조",
        "concepts": [
            ("메모리 4개 영역", [
                "코드(Code) 영역: 실행할 프로그램 코드 저장 (읽기 전용)",
                "데이터(Data) 영역: 전역 변수, static 변수 저장",
                "힙(Heap) 영역: 동적 할당 메모리 (↑ 위로 성장)",
                "스택(Stack) 영역: 지역 변수, 매개변수 (↓ 아래로 성장)"
            ]),
            ("64비트 주소 체계", [
                "모든 포인터 크기: 8바이트 (64비트)",
                "주소 범위: 0x0000000000000000 ~ 0xFFFFFFFFFFFFFFFF",
                "사용자 공간과 커널 공간으로 분리"
            ])
        ],
        "code": r'''#include <stdio.h>
int global_var = 100;
void my_function() { printf("Hello!\n"); }
int main() {
    int local_var = 200;
    int *heap_var = (int*)malloc(sizeof(int));
    *heap_var = 300;
    printf("코드 영역(함수):  %p\n", (void*)my_function);
    printf("데이터 영역(전역): %p\n", (void*)&global_var);
    printf("힙 영역(동적):    %p\n", (void*)heap_var);
    printf("스택 영역(지역):  %p\n", (void*)&local_var);
    free(heap_var);
    return 0;
}''',
        "memory": [
            ("0x00401130", "my_function", "(함수코드)", "function", "func"),
            ("0x00604020", "global_var", "100", "int", "int"),
            ("0x00B71A40", "*heap_var", "300", "int (힙)", "heap"),
            ("0x7FFE0028", "local_var", "200", "int (스택)", "int"),
            ("0x7FFE0030", "heap_var", "0x00B71A40", "int* (스택)", "ptr"),
        ],
        "arrows": "heap_var(0x7FFE0030) → *heap_var(0x00B71A40) : 스택의 포인터가 힙 메모리를 가리킴",
        "keypoints": [
            "코드 영역은 읽기 전용으로 프로그램 실행 코드 저장",
            "데이터 영역은 전역 변수와 static 변수 저장",
            "힙 영역은 동적 할당 메모리 (위로 성장 ↑)",
            "스택 영역은 지역 변수 저장 (아래로 성장 ↓)",
            "64비트 시스템에서 포인터 크기는 항상 8바이트"
        ]
    },
    {
        "id": 2,
        "title": "변수와 메모리",
        "concepts": [
            ("자료형별 크기", [
                "char: 1바이트 (문자 1개)",
                "short: 2바이트",
                "int: 4바이트 (정수)",
                "long long: 8바이트 (큰 정수)",
                "float: 4바이트 (단정밀도 실수)",
                "double: 8바이트 (배정밀도 실수)",
                "포인터: 8바이트 (64비트 시스템)"
            ]),
            ("메모리 할당 과정", [
                "변수 선언 → 컴파일러가 메모리 공간 확보",
                "자료형에 따라 할당 크기 결정",
                "sizeof 연산자로 크기 확인 가능"
            ])
        ],
        "code": r'''#include <stdio.h>
int main() {
    char c = 'A';
    short s = 100;
    int i = 42;
    long long ll = 123456789LL;
    float f = 3.14f;
    double d = 2.71828;
    printf("char:      %zu byte, addr: %p\n", sizeof(c), &c);
    printf("short:     %zu byte, addr: %p\n", sizeof(s), &s);
    printf("int:       %zu byte, addr: %p\n", sizeof(i), &i);
    printf("long long: %zu byte, addr: %p\n", sizeof(ll), &ll);
    printf("float:     %zu byte, addr: %p\n", sizeof(f), &f);
    printf("double:    %zu byte, addr: %p\n", sizeof(d), &d);
    return 0;
}''',
        "memory": [
            ("0x7FFE0047", "c", "'A' (0x41)", "char (1B)", "char"),
            ("0x7FFE0044", "s", "100", "short (2B)", "int"),
            ("0x7FFE0040", "i", "42", "int (4B)", "int"),
            ("0x7FFE0038", "ll", "123456789", "long long (8B)", "int"),
            ("0x7FFE0034", "f", "3.14", "float (4B)", "int"),
            ("0x7FFE0028", "d", "2.71828", "double (8B)", "int"),
        ],
        "arrows": "스택에 연속 배치: 높은 주소(c) → 낮은 주소(d) 순서로 할당",
        "keypoints": [
            "자료형은 메모리 접근 크기를 결정한다",
            "변수 선언 = 메모리 공간 확보",
            "sizeof 연산자로 바이트 크기 확인 가능",
            "64비트 시스템에서 모든 포인터는 8바이트"
        ]
    },
    {
        "id": 3,
        "title": "포인터의 개념",
        "concepts": [
            ("포인터란?", [
                "포인터 = 메모리 주소를 저장하는 변수",
                "& 연산자: 변수의 주소를 얻는 연산자",
                "* 연산자: 주소가 가리키는 값에 접근 (역참조)",
            ]),
            ("64비트 포인터", [
                "모든 포인터는 8바이트 크기",
                "자료형에 상관없이 주소 저장 크기는 동일",
                "포인터 자체도 메모리에 저장됨 (주소를 가짐)"
            ])
        ],
        "code": r'''#include <stdio.h>
int main() {
    int num = 42;
    int *ptr = &num;
    printf("num의 값:    %d\n", num);
    printf("num의 주소:  %p\n", (void*)&num);
    printf("ptr의 값:    %p\n", (void*)ptr);
    printf("ptr의 주소:  %p\n", (void*)&ptr);
    printf("*ptr의 값:   %d\n", *ptr);
    printf("포인터 크기: %zu\n", sizeof(ptr));
    return 0;
}''',
        "memory": [
            ("0x7FFE0038", "num", "42", "int (4B)", "int"),
            ("0x7FFE0040", "ptr", "0x7FFE0038", "int* (8B)", "ptr"),
        ],
        "arrows": "ptr(0x7FFE0040) → num(0x7FFE0038) : ptr이 num의 주소를 저장",
        "keypoints": [
            "포인터는 메모리 주소를 저장하는 변수",
            "& 연산자로 변수의 주소를 얻음",
            "* 연산자로 주소가 가리키는 값에 접근 (역참조)",
            "64비트 시스템에서 포인터는 항상 8바이트"
        ]
    },
    {
        "id": 4,
        "title": "포인터 선언과 초기화",
        "concepts": [
            ("포인터 선언 문법", [
                "int *pa = &a;     // int형 포인터",
                "char *pb = &b;    // char형 포인터",
                "double *pc = &c;  // double형 포인터",
            ]),
            ("NULL 초기화", [
                "int *pn = NULL;  // 안전한 초기화",
                "NULL = 아무것도 가리키지 않음을 명시",
                "사용 전 NULL 체크로 안전성 확보"
            ]),
            ("포인터 크기", [
                "모든 포인터는 자료형에 무관하게 8바이트 (64비트)",
                "int*, char*, double* 모두 같은 크기"
            ])
        ],
        "code": r'''#include <stdio.h>
int main() {
    int a = 10;
    char b = 'X';
    double c = 3.14;
    int *pa = &a;
    char *pb = &b;
    double *pc = &c;
    int *pn = NULL;
    printf("int* pa = %p -> %d\n", pa, *pa);
    printf("char* pb = %p -> %c\n", pb, *pb);
    printf("double* pc = %p -> %.2f\n", pc, *pc);
    printf("NULL ptr: %p\n", pn);
    return 0;
}''',
        "memory": [
            ("0x7FFE0020", "a", "10", "int (4B)", "int"),
            ("0x7FFE0027", "b", "'X' (0x58)", "char (1B)", "char"),
            ("0x7FFE0028", "c", "3.14", "double (8B)", "int"),
            ("0x7FFE0038", "pa", "0x7FFE0020", "int* (8B)", "ptr"),
            ("0x7FFE0040", "pb", "0x7FFE0027", "char* (8B)", "ptr"),
            ("0x7FFE0048", "pc", "0x7FFE0028", "double* (8B)", "ptr"),
            ("0x7FFE0050", "pn", "NULL (0x0)", "int* (8B)", "ptr"),
        ],
        "arrows": "pa → a, pb → b, pc → c : 각 포인터가 해당 변수를 가리킴 / pn = NULL (미가리킴)",
        "keypoints": [
            "자료형과 일치하는 포인터 선언 필수",
            "NULL 초기화로 안전성 확보",
            "포인터 크기는 자료형에 무관하게 8바이트 (64비트)",
            "사용 전 NULL 체크 습관화"
        ]
    },
    {
        "id": 5,
        "title": "역참조 연산자(*)",
        "concepts": [
            ("역참조(Dereference)란?", [
                "* 연산자로 포인터가 가리키는 값에 접근",
                "*p 로 값 읽기 (Read)",
                "*p = 99; 로 값 쓰기 (Write)",
            ]),
            ("포인터를 통한 간접 변경", [
                "포인터로 값 변경 시 원본 변수도 변경됨",
                "원본 변수 변경 시 포인터로 읽은 값도 변경됨",
                "p = ... (주소 변경) vs *p = ... (값 변경) 구분 중요"
            ])
        ],
        "code": r'''#include <stdio.h>
int main() {
    int a = 10;
    int *p = &a;
    printf("변경 전: a = %d, *p = %d\n", a, *p);
    *p = 99;
    printf("변경 후: a = %d, *p = %d\n", a, *p);
    a = 50;
    printf("재변경:  a = %d, *p = %d\n", a, *p);
    return 0;
}''',
        "memory": [
            ("0x7FFE0038", "a", "10 → 99 → 50", "int (4B)", "int"),
            ("0x7FFE0040", "p", "0x7FFE0038", "int* (8B)", "ptr"),
        ],
        "arrows": "p → a : *p = 99 실행 시 a의 값이 10에서 99로 변경됨",
        "keypoints": [
            "*p로 읽기와 쓰기 모두 가능",
            "포인터로 값 변경 시 원본 변수도 변경",
            "p = (주소 변경) vs *p = (값 변경) 구분 중요",
            "역참조는 포인터의 핵심 기능"
        ]
    },
    {
        "id": 6,
        "title": "포인터와 자료형의 관계",
        "concepts": [
            ("자료형이 결정하는 것", [
                "포인터 자료형 = 메모리에서 몇 바이트를 읽을지 결정",
                "int* → 4바이트 읽기/쓰기",
                "char* → 1바이트 읽기/쓰기",
                "double* → 8바이트 읽기/쓰기",
            ]),
            ("void 포인터 (void*)", [
                "범용 포인터: 어떤 주소든 저장 가능",
                "직접 역참조 불가 → 캐스팅 필요",
                "malloc 반환 타입이 void*"
            ])
        ],
        "code": r'''#include <stdio.h>
int main() {
    int num = 0x41424344;
    int *ip = &num;
    char *cp = (char*)&num;
    printf("int*로 읽기: 0x%X\n", *ip);
    printf("char*로 읽기:\n");
    for(int i = 0; i < 4; i++)
        printf("  바이트[%d] = 0x%02X ('%c')\n",
               i, (unsigned char)cp[i], cp[i]);
    void *vp = &num;
    printf("void* 주소: %p\n", vp);
    printf("void* → int*: %d\n", *(int*)vp);
    return 0;
}''',
        "memory": [
            ("0x7FFE0038", "num", "0x41424344", "int (4B)", "int"),
            ("0x7FFE0038+0", "바이트[0]", "0x44 ('D')", "char (1B)", "char"),
            ("0x7FFE0038+1", "바이트[1]", "0x43 ('C')", "char (1B)", "char"),
            ("0x7FFE0038+2", "바이트[2]", "0x42 ('B')", "char (1B)", "char"),
            ("0x7FFE0038+3", "바이트[3]", "0x41 ('A')", "char (1B)", "char"),
        ],
        "arrows": "int*는 4바이트 전체를 한 번에 읽음 / char*는 1바이트씩 개별 접근",
        "keypoints": [
            "포인터 자료형은 메모리 읽기/쓰기 단위를 결정",
            "같은 주소라도 자료형에 따라 다른 값으로 해석",
            "void*는 범용 포인터 (역참조 불가, 캐스팅 필요)",
            "리틀 엔디안: 낮은 주소에 하위 바이트 저장"
        ]
    },
    {
        "id": 7,
        "title": "포인터 연산",
        "concepts": [
            ("포인터 산술 연산", [
                "포인터 + n = 주소 + (n × sizeof(자료형))",
                "int* + 1 = 주소 + 4바이트",
                "char* + 1 = 주소 + 1바이트",
                "double* + 1 = 주소 + 8바이트",
            ]),
            ("포인터 뺄셈", [
                "포인터 - 포인터 = 두 주소 사이의 요소 개수",
                "같은 배열 내에서만 의미 있음",
                "결과 타입: ptrdiff_t"
            ])
        ],
        "code": r'''#include <stdio.h>
int main() {
    int arr[] = {10, 20, 30, 40, 50};
    int *p = arr;
    printf("p   = %p, *p   = %d\n", p, *p);
    printf("p+1 = %p, *(p+1)= %d\n", p+1, *(p+1));
    printf("p+2 = %p, *(p+2)= %d\n", p+2, *(p+2));
    char *cp = (char*)arr;
    printf("\nchar* cp   = %p\n", cp);
    printf("char* cp+1 = %p (1바이트 이동)\n", cp+1);
    int *end = &arr[4];
    printf("\n포인터 차이: end - p = %td\n", end - p);
    return 0;
}''',
        "memory": [
            ("0x7FFE0020", "arr[0]", "10", "int (4B)", "arr"),
            ("0x7FFE0024", "arr[1]", "20", "int (4B)", "arr"),
            ("0x7FFE0028", "arr[2]", "30", "int (4B)", "arr"),
            ("0x7FFE002C", "arr[3]", "40", "int (4B)", "arr"),
            ("0x7FFE0030", "arr[4]", "50", "int (4B)", "arr"),
            ("0x7FFE0038", "p", "0x7FFE0020", "int* (8B)", "ptr"),
        ],
        "arrows": "p → arr[0], p+1 → arr[1], p+2 → arr[2] : int* +1 = 4바이트씩 이동",
        "keypoints": [
            "int* + 1 = 4바이트 이동",
            "char* + 1 = 1바이트 이동",
            "포인터 뺄셈 = 요소 개수 차이 반환",
            "포인터 연산은 자료형 크기 기반"
        ]
    },
    {
        "id": 8,
        "title": "배열과 메모리",
        "concepts": [
            ("배열의 메모리 구조", [
                "배열 = 연속된 메모리 공간에 같은 타입 저장",
                "배열명 = 첫 번째 요소의 주소 (상수 포인터)",
                "arr[i] == *(arr + i) : 동일한 접근 방법",
            ]),
            ("배열 크기 계산", [
                "sizeof(arr) = 배열 전체 크기 (바이트)",
                "sizeof(arr[0]) = 요소 하나의 크기",
                "배열 길이 = sizeof(arr) / sizeof(arr[0])"
            ])
        ],
        "code": r'''#include <stdio.h>
int main() {
    int arr[5] = {10, 20, 30, 40, 50};
    printf("배열 시작 주소: %p\n", (void*)arr);
    for(int i = 0; i < 5; i++) {
        printf("arr[%d]: 주소=%p, 값=%d\n",
               i, (void*)&arr[i], arr[i]);
    }
    printf("\nsizeof(arr) = %zu\n", sizeof(arr));
    printf("sizeof(arr[0]) = %zu\n", sizeof(arr[0]));
    printf("배열 길이 = %zu\n",
           sizeof(arr)/sizeof(arr[0]));
    return 0;
}''',
        "memory": [
            ("0x7FFE0020", "arr[0]", "10", "int (4B)", "arr"),
            ("0x7FFE0024", "arr[1]", "20", "int (4B)", "arr"),
            ("0x7FFE0028", "arr[2]", "30", "int (4B)", "arr"),
            ("0x7FFE002C", "arr[3]", "40", "int (4B)", "arr"),
            ("0x7FFE0030", "arr[4]", "50", "int (4B)", "arr"),
        ],
        "arrows": "arr = &arr[0] = 0x7FFE0020 : 배열명은 첫 번째 요소의 주소 (연속 배치, 4B 간격)",
        "keypoints": [
            "배열은 메모리에 연속 배치됨",
            "배열명은 상수 포인터 (주소 변경 불가)",
            "sizeof(배열) = 전체 크기 (바이트)",
            "배열 길이 = sizeof(arr) / sizeof(arr[0])"
        ]
    },
    {
        "id": 9,
        "title": "배열과 포인터의 관계",
        "concepts": [
            ("배열과 포인터의 등가 관계", [
                "arr[i] == *(arr+i) == p[i] == *(p+i)",
                "4가지 표현 모두 동일한 결과",
                "포인터로 배열의 모든 요소 접근 가능",
            ]),
            ("배열명 vs 포인터의 차이", [
                "배열명: 상수 → 주소 변경 불가 (arr++ 불가)",
                "포인터: 변수 → 주소 변경 가능 (p++ 가능)",
                "sizeof 차이: sizeof(arr)=전체, sizeof(p)=8"
            ])
        ],
        "code": r'''#include <stdio.h>
int main() {
    int arr[5] = {10, 20, 30, 40, 50};
    int *p = arr;
    printf("=== 배열 첨자 vs 포인터 연산 ===\n");
    for(int i = 0; i < 5; i++) {
        printf("arr[%d]=%d  *(arr+%d)=%d  "
               "p[%d]=%d  *(p+%d)=%d\n",
               i, arr[i], i, *(arr+i),
               i, p[i], i, *(p+i));
    }
    printf("\n=== 포인터 이동으로 순회 ===\n");
    int *q = arr;
    while(q < arr + 5) {
        printf("%d ", *q);
        q++;
    }
    printf("\n");
    return 0;
}''',
        "memory": [
            ("0x7FFE0020", "arr[0]", "10", "int (4B)", "arr"),
            ("0x7FFE0024", "arr[1]", "20", "int (4B)", "arr"),
            ("0x7FFE0028", "arr[2]", "30", "int (4B)", "arr"),
            ("0x7FFE002C", "arr[3]", "40", "int (4B)", "arr"),
            ("0x7FFE0030", "arr[4]", "50", "int (4B)", "arr"),
            ("0x7FFE0038", "p", "0x7FFE0020", "int* (8B)", "ptr"),
        ],
        "arrows": "p → arr[0] : arr[i] == *(arr+i) == p[i] == *(p+i) (4가지 동일 접근)",
        "keypoints": [
            "arr[i] == *(arr+i) == p[i] == *(p+i)",
            "배열명 = 주소 변경 불가 (상수 포인터)",
            "포인터 = 주소 변경 가능 (변수)",
            "sizeof(arr) = 전체 크기, sizeof(p) = 8바이트"
        ]
    },
    {
        "id": 10,
        "title": "문자열과 포인터",
        "concepts": [
            ("char 배열 vs char 포인터", [
                "char str1[] = \"Hello\";  → 스택에 복사 (수정 가능)",
                "char *str2 = \"World\";   → 데이터 영역 가리킴 (수정 불가)",
                "sizeof(str1) = 6 (문자 + '\\0')",
                "sizeof(str2) = 8 (포인터 크기)",
            ]),
            ("문자열 순회", [
                "포인터로 문자열 순회: while(*p) { ... p++; }",
                "'\\0' (NULL 문자)을 만나면 종료",
                "문자열 = char 배열 + '\\0' 종단 문자"
            ])
        ],
        "code": r'''#include <stdio.h>
#include <string.h>
int main() {
    char str1[] = "Hello";
    char *str2 = "World";
    printf("str1: %s (주소: %p, 크기: %zu)\n",
           str1, str1, sizeof(str1));
    printf("str2: %s (주소: %p, 크기: %zu)\n",
           str2, str2, sizeof(str2));
    str1[0] = 'h';
    printf("str1 수정 후: %s\n", str1);
    // str2[0] = 'w';  // SEGFAULT!
    char *p = str1;
    printf("포인터로 순회: ");
    while(*p) { printf("%c ", *p); p++; }
    printf("\n");
    return 0;
}''',
        "memory": [
            ("0x00404000", "\"World\"", "W,o,r,l,d,\\0", "데이터 영역 (읽기전용)", "func"),
            ("0x7FFE0030", "str1[0..5]", "H,e,l,l,o,\\0", "char[6] (스택, 수정가능)", "char"),
            ("0x7FFE0040", "str2", "0x00404000", "char* (8B)", "ptr"),
        ],
        "arrows": "str2(0x7FFE0040) → \"World\"(0x00404000) : 데이터 영역의 문자열 리터럴 가리킴",
        "keypoints": [
            "char[]은 스택에 복사 → 수정 가능",
            "char*는 데이터 영역 → 수정 불가 (읽기 전용)",
            "sizeof 차이: char[] = 문자열 크기, char* = 8바이트",
            "문자열 끝은 항상 '\\0' (NULL 문자)"
        ]
    },
    {
        "id": 11,
        "title": "Call by Value vs Call by Reference",
        "concepts": [
            ("Call by Value (값에 의한 호출)", [
                "함수에 값의 복사본을 전달",
                "함수 내부에서 변경해도 원본은 불변",
                "swap_value(x, y) → x, y 변경 안 됨",
            ]),
            ("Call by Reference (참조에 의한 호출)", [
                "함수에 변수의 주소를 전달",
                "포인터를 통해 원본 값을 직접 변경",
                "swap_ref(&x, &y) → x, y 실제 변경됨"
            ])
        ],
        "code": r'''#include <stdio.h>
void swap_value(int a, int b) {
    int temp = a; a = b; b = temp;
    printf("  함수 내부: a=%d, b=%d\n", a, b);
}
void swap_ref(int *a, int *b) {
    int temp = *a; *a = *b; *b = temp;
    printf("  함수 내부: *a=%d, *b=%d\n", *a, *b);
}
int main() {
    int x = 10, y = 20;
    printf("Call by Value 전: x=%d, y=%d\n", x, y);
    swap_value(x, y);
    printf("Call by Value 후: x=%d, y=%d\n\n", x, y);
    printf("Call by Reference 전: x=%d, y=%d\n", x, y);
    swap_ref(&x, &y);
    printf("Call by Reference 후: x=%d, y=%d\n", x, y);
    return 0;
}''',
        "memory": [
            ("0x7FFE0020", "x (main)", "10 → 20", "int (4B)", "int"),
            ("0x7FFE0024", "y (main)", "20 → 10", "int (4B)", "int"),
            ("0x7FFE0010", "a (swap_val)", "10 (복사본)", "int (4B)", "char"),
            ("0x7FFE0014", "b (swap_val)", "20 (복사본)", "int (4B)", "char"),
            ("0x7FFE0008", "*a (swap_ref)", "→ x 가리킴", "int* (8B)", "ptr"),
            ("0x7FFE0010", "*b (swap_ref)", "→ y 가리킴", "int* (8B)", "ptr"),
        ],
        "arrows": "swap_value: a,b는 x,y의 복사본 / swap_ref: *a→x, *b→y (원본 직접 접근)",
        "keypoints": [
            "Call by Value = 복사본 전달 → 원본 불변",
            "Call by Reference = 주소 전달 → 원본 변경 가능",
            "C언어에서 참조 전달은 포인터로 구현",
            "swap 함수는 포인터 활용의 대표적 예시"
        ]
    },
    {
        "id": 12,
        "title": "포인터 매개변수 활용",
        "concepts": [
            ("포인터 매개변수의 활용", [
                "여러 값을 동시에 반환 (출력 매개변수)",
                "배열을 함수에 전달 (포인터로 자동 변환)",
                "구조체 포인터로 큰 데이터 효율적 전달",
            ]),
            ("배열 매개변수", [
                "배열은 항상 포인터로 전달됨",
                "함수 내에서 배열 크기 알 수 없음 → 크기 별도 전달",
                "void func(int *arr, int size) 패턴"
            ])
        ],
        "code": r'''#include <stdio.h>
void getMinMax(int *arr, int size,
               int *min, int *max) {
    *min = *max = arr[0];
    for(int i = 1; i < size; i++) {
        if(arr[i] < *min) *min = arr[i];
        if(arr[i] > *max) *max = arr[i];
    }
}
void reverseArray(int *arr, int size) {
    for(int i = 0; i < size/2; i++) {
        int temp = arr[i];
        arr[i] = arr[size-1-i];
        arr[size-1-i] = temp;
    }
}
int main() {
    int data[] = {38, 7, 91, 23, 55};
    int min, max;
    getMinMax(data, 5, &min, &max);
    printf("최솟값: %d, 최댓값: %d\n", min, max);
    reverseArray(data, 5);
    printf("뒤집기: ");
    for(int i=0;i<5;i++) printf("%d ",data[i]);
    return 0;
}''',
        "memory": [
            ("0x7FFE0030", "data[0..4]", "{38,7,91,23,55}", "int[5] (20B)", "arr"),
            ("0x7FFE0048", "min", "7", "int (4B)", "int"),
            ("0x7FFE004C", "max", "91", "int (4B)", "int"),
            ("0x7FFE0010", "arr (매개변수)", "→ data 가리킴", "int* (8B)", "ptr"),
            ("0x7FFE0018", "min (매개변수)", "→ min 가리킴", "int* (8B)", "ptr"),
            ("0x7FFE0020", "max (매개변수)", "→ max 가리킴", "int* (8B)", "ptr"),
        ],
        "arrows": "함수의 포인터 매개변수들이 main의 변수들을 직접 가리켜 값을 변경",
        "keypoints": [
            "포인터 매개변수로 다중 값 반환 가능",
            "배열은 함수에 항상 포인터로 전달",
            "배열 크기는 별도 매개변수로 전달 필요",
            "출력 매개변수 패턴으로 효율적 데이터 반환"
        ]
    },
    {
        "id": 13,
        "title": "함수의 포인터 리턴",
        "concepts": [
            ("안전한 포인터 반환", [
                "동적 할당(malloc) 메모리 주소 → 안전",
                "static 변수 주소 → 안전 (프로그램 종료까지 유지)",
                "호출자가 free() 책임",
            ]),
            ("위험한 포인터 반환 (댕글링 포인터)", [
                "지역 변수 주소 반환 → 위험! (함수 종료 시 소멸)",
                "댕글링 포인터: 이미 해제된 메모리를 가리키는 포인터",
                "정의되지 않은 동작(Undefined Behavior) 발생"
            ])
        ],
        "code": r'''#include <stdio.h>
#include <stdlib.h>
int* createArray(int size) {
    int *arr = (int*)malloc(sizeof(int) * size);
    for(int i = 0; i < size; i++)
        arr[i] = (i+1)*10;
    return arr;
}
int* getCounter() {
    static int count = 0;
    count++;
    return &count;
}
// int* dangerous() {
//     int x=42; return &x; // 위험!
// }
int main() {
    int *arr = createArray(3);
    printf("동적 배열: %d %d %d\n",
           arr[0], arr[1], arr[2]);
    int *c1 = getCounter();
    int *c2 = getCounter();
    printf("카운터: %d\n", *c2);
    free(arr);
    return 0;
}''',
        "memory": [
            ("0x00B71A40", "arr[0] (힙)", "10", "int (동적)", "heap"),
            ("0x00B71A44", "arr[1] (힙)", "20", "int (동적)", "heap"),
            ("0x00B71A48", "arr[2] (힙)", "30", "int (동적)", "heap"),
            ("0x00604030", "count (static)", "2", "int (데이터영역)", "int"),
            ("0x7FFE0038", "arr (스택)", "0x00B71A40", "int* (8B)", "ptr"),
        ],
        "arrows": "arr(스택) → 힙 영역 동적 배열 / getCounter → static 변수 (데이터 영역, 안전)",
        "keypoints": [
            "동적 할당(malloc)이나 static 변수 주소만 안전하게 반환 가능",
            "지역 변수 주소 반환 = 댕글링 포인터 (위험!)",
            "동적 할당 메모리는 반드시 free()로 해제",
            "static 변수는 프로그램 종료까지 유지됨"
        ]
    },
    {
        "id": 14,
        "title": "const와 포인터",
        "concepts": [
            ("const 포인터의 3가지 형태", [
                "const int *p  → 값 변경 불가, 주소 변경 가능",
                "int * const p → 값 변경 가능, 주소 변경 불가",
                "const int * const p → 둘 다 변경 불가",
            ]),
            ("읽는 방법", [
                "* 기준 왼쪽에 const → 값 보호",
                "* 기준 오른쪽에 const → 주소 고정",
                "API 설계 시 const 활용: 읽기 전용 보장"
            ])
        ],
        "code": r'''#include <stdio.h>
int main() {
    int a = 10, b = 20;
    const int *p1 = &a;
    // *p1 = 30;  // 에러! 값 변경 불가
    p1 = &b;      // OK! 주소 변경 가능
    printf("const int *p1: %d\n", *p1);
    int * const p2 = &a;
    *p2 = 30;     // OK! 값 변경 가능
    // p2 = &b;   // 에러! 주소 변경 불가
    printf("int * const p2: %d\n", *p2);
    const int * const p3 = &a;
    // *p3 = 40;  // 에러!
    // p3 = &b;   // 에러!
    printf("const int * const p3: %d\n", *p3);
    return 0;
}''',
        "memory": [
            ("0x7FFE0020", "a", "30 (p2가 변경)", "int (4B)", "int"),
            ("0x7FFE0024", "b", "20", "int (4B)", "int"),
            ("0x7FFE0028", "p1", "→ b (주소변경O)", "const int*", "ptr"),
            ("0x7FFE0030", "p2", "→ a (주소고정)", "int* const", "ptr"),
            ("0x7FFE0038", "p3", "→ a (둘다고정)", "const int* const", "ptr"),
        ],
        "arrows": "p1: 값변경X 주소변경O / p2: 값변경O 주소변경X / p3: 값변경X 주소변경X",
        "keypoints": [
            "const int* p = 값 보호 (읽기 전용 접근)",
            "int* const p = 주소 고정 (다른 변수 가리킬 수 없음)",
            "const int* const p = 값 보호 + 주소 고정",
            "읽기 방법: * 기준 왼쪽=값보호, 오른쪽=주소고정"
        ]
    },
    {
        "id": 15,
        "title": "동적 메모리 할당 기초",
        "concepts": [
            ("동적 메모리 할당 함수", [
                "malloc(size): size 바이트 할당 (초기화 안됨)",
                "calloc(n, size): n×size 바이트 할당 (0 초기화)",
                "realloc(ptr, size): 기존 메모리 크기 변경",
                "free(ptr): 할당된 메모리 해제",
            ]),
            ("힙 영역 활용", [
                "프로그래머가 직접 할당/해제 관리",
                "런타임에 크기 결정 가능",
                "free 하지 않으면 메모리 누수 발생"
            ])
        ],
        "code": r'''#include <stdio.h>
#include <stdlib.h>
int main() {
    int *p1 = (int*)malloc(sizeof(int) * 3);
    p1[0]=10; p1[1]=20; p1[2]=30;
    int *p2 = (int*)calloc(3, sizeof(int));
    printf("malloc: %d %d %d\n",
           p1[0], p1[1], p1[2]);
    printf("calloc: %d %d %d\n",
           p2[0], p2[1], p2[2]);
    p1 = (int*)realloc(p1, sizeof(int)*5);
    p1[3]=40; p1[4]=50;
    printf("realloc: ");
    for(int i=0;i<5;i++) printf("%d ",p1[i]);
    printf("\n");
    free(p1); free(p2);
    return 0;
}''',
        "memory": [
            ("0x7FFE0038", "p1 (스택)", "0x00B71A40", "int* (8B)", "ptr"),
            ("0x7FFE0040", "p2 (스택)", "0x00B71A60", "int* (8B)", "ptr"),
            ("0x00B71A40", "p1[0] (힙)", "10", "int (malloc)", "heap"),
            ("0x00B71A44", "p1[1] (힙)", "20", "int (malloc)", "heap"),
            ("0x00B71A48", "p1[2] (힙)", "30", "int (malloc)", "heap"),
            ("0x00B71A60", "p2[0] (힙)", "0", "int (calloc)", "heap"),
        ],
        "arrows": "p1(스택) → malloc 영역(힙) / p2(스택) → calloc 영역(힙, 0초기화)",
        "keypoints": [
            "malloc = 초기화 안 됨 (쓰레기 값)",
            "calloc = 0으로 초기화",
            "realloc = 기존 메모리 크기 변경",
            "반드시 free()로 해제 (메모리 누수 방지)"
        ]
    },
    {
        "id": 16,
        "title": "동적 메모리와 포인터 활용",
        "concepts": [
            ("동적 2D 배열", [
                "이중 포인터 + 이중 malloc으로 구현",
                "int **matrix = malloc(rows × sizeof(int*))",
                "matrix[i] = malloc(cols × sizeof(int))",
            ]),
            ("메모리 관리 주의사항", [
                "역순 free 필수: 행 먼저 → 포인터 배열 나중",
                "free 후 포인터를 NULL로 설정 (안전)",
                "메모리 누수 = free 하지 않은 동적 메모리"
            ])
        ],
        "code": r'''#include <stdio.h>
#include <stdlib.h>
#include <string.h>
int main() {
    int rows=2, cols=3;
    int **matrix =
        (int**)malloc(sizeof(int*)*rows);
    for(int i=0; i<rows; i++)
        matrix[i] =
            (int*)malloc(sizeof(int)*cols);
    int val=1;
    for(int i=0;i<rows;i++)
        for(int j=0;j<cols;j++)
            matrix[i][j] = val++;
    printf("2D 행렬:\n");
    for(int i=0;i<rows;i++){
        for(int j=0;j<cols;j++)
            printf("%d ", matrix[i][j]);
        printf("\n");
    }
    for(int i=0;i<rows;i++) free(matrix[i]);
    free(matrix);
    return 0;
}''',
        "memory": [
            ("0x7FFE0038", "matrix (스택)", "0x00B71A40", "int** (8B)", "ptr"),
            ("0x00B71A40", "matrix[0] (힙)", "0x00B71A60", "int* (8B)", "heap"),
            ("0x00B71A48", "matrix[1] (힙)", "0x00B71A80", "int* (8B)", "heap"),
            ("0x00B71A60", "[0][0],[0][1],[0][2]", "1, 2, 3", "int[3] (힙)", "heap"),
            ("0x00B71A80", "[1][0],[1][1],[1][2]", "4, 5, 6", "int[3] (힙)", "heap"),
        ],
        "arrows": "matrix(스택) → int* 배열(힙) → 각 행 int 배열(힙) : 3단계 참조",
        "keypoints": [
            "2D 동적 배열 = 이중 포인터 + 이중 malloc",
            "역순 free 필수 (행 먼저 → 포인터 배열 나중)",
            "메모리 누수 주의: 모든 malloc에 대응하는 free",
            "free 후 포인터를 NULL로 설정하면 안전"
        ]
    },
    {
        "id": 17,
        "title": "이중 포인터 (더블 포인터)",
        "concepts": [
            ("이중 포인터란?", [
                "int **pp = 포인터의 주소를 저장하는 변수",
                "3단계 참조: pp → p → 변수",
                "**pp로 원본 값에 접근",
            ]),
            ("이중 포인터의 활용", [
                "함수에서 포인터 자체를 변경할 때 필수",
                "동적 할당 함수: void alloc(int **pp)",
                "2D 배열의 동적 할당"
            ])
        ],
        "code": r'''#include <stdio.h>
#include <stdlib.h>
void allocate(int **pp, int value) {
    *pp = (int*)malloc(sizeof(int));
    **pp = value;
}
int main() {
    int a = 10;
    int *p = &a;
    int **pp = &p;
    printf("a = %d\n", a);
    printf("*p = %d\n", *p);
    printf("**pp = %d\n", **pp);
    printf("\n주소 관계:\n");
    printf("&a  = %p\n", &a);
    printf("p   = %p\n", p);
    printf("&p  = %p\n", &p);
    printf("pp  = %p\n", pp);
    int *dynamic = NULL;
    allocate(&dynamic, 42);
    printf("\n동적 할당: %d\n", *dynamic);
    free(dynamic);
    return 0;
}''',
        "memory": [
            ("0x7FFE0020", "a", "10", "int (4B)", "int"),
            ("0x7FFE0028", "p", "0x7FFE0020", "int* (8B)", "ptr"),
            ("0x7FFE0030", "pp", "0x7FFE0028", "int** (8B)", "ptr"),
        ],
        "arrows": "pp(0x7FFE0030) → p(0x7FFE0028) → a(0x7FFE0020) : 3단계 참조 체인",
        "keypoints": [
            "이중 포인터 = 포인터의 포인터",
            "**pp로 원본 값에 접근 (3단계 참조)",
            "함수에서 포인터 자체를 변경할 때 필수",
            "동적 할당 함수 패턴: void alloc(int **pp)"
        ]
    },
    {
        "id": 18,
        "title": "포인터 배열과 배열 포인터",
        "concepts": [
            ("포인터 배열: int *arr[3]", [
                "포인터 3개를 담는 배열",
                "sizeof = 3 × 8 = 24바이트 (64비트)",
                "각 요소가 서로 다른 변수를 가리킬 수 있음",
            ]),
            ("배열 포인터: int (*arr)[3]", [
                "크기 3인 int 배열을 가리키는 포인터",
                "sizeof = 8바이트 (포인터 1개)",
                "2D 배열 접근에 활용",
                "+1 시 배열 크기(12B)만큼 이동"
            ])
        ],
        "code": r'''#include <stdio.h>
int main() {
    int a=10, b=20, c=30;
    int *pArr[3] = {&a, &b, &c};
    printf("=== 포인터 배열 int *pArr[3] ===\n");
    for(int i=0;i<3;i++)
        printf("pArr[%d] = %p -> %d\n",
               i, pArr[i], *pArr[i]);
    int matrix[2][3] = {{1,2,3},{4,5,6}};
    int (*rowPtr)[3] = matrix;
    printf("\n=== 배열 포인터 int (*rowPtr)[3]"
           " ===\n");
    for(int i=0;i<2;i++)
        for(int j=0;j<3;j++)
            printf("rowPtr[%d][%d] = %d\n",
                   i, j, rowPtr[i][j]);
    printf("\nsizeof(pArr)=%zu, "
           "sizeof(rowPtr)=%zu\n",
           sizeof(pArr), sizeof(rowPtr));
    return 0;
}''',
        "memory": [
            ("0x7FFE0020", "a", "10", "int (4B)", "int"),
            ("0x7FFE0024", "b", "20", "int (4B)", "int"),
            ("0x7FFE0028", "c", "30", "int (4B)", "int"),
            ("0x7FFE0030", "pArr[0]", "→ a", "int* (8B)", "ptr"),
            ("0x7FFE0038", "pArr[1]", "→ b", "int* (8B)", "ptr"),
            ("0x7FFE0040", "pArr[2]", "→ c", "int* (8B)", "ptr"),
            ("0x7FFE0050", "rowPtr", "→ matrix", "int(*)[3] (8B)", "ptr"),
        ],
        "arrows": "포인터 배열: pArr[0]→a, pArr[1]→b, pArr[2]→c / 배열 포인터: rowPtr→matrix 전체",
        "keypoints": [
            "int *arr[3] = 포인터 3개의 배열 (24B)",
            "int (*arr)[3] = 크기3 배열을 가리키는 포인터 (8B)",
            "괄호 유무가 의미를 완전히 바꿈",
            "배열 포인터는 2D 배열 접근에 활용"
        ]
    },
    {
        "id": 19,
        "title": "함수 포인터",
        "concepts": [
            ("함수 포인터란?", [
                "함수도 메모리에 주소를 가짐 (코드 영역)",
                "함수명 = 함수의 시작 주소",
                "int (*fp)(int, int) = add; 형태로 선언",
            ]),
            ("활용 패턴", [
                "콜백(Callback) 패턴 구현",
                "함수 포인터 배열로 전략 패턴",
                "typedef로 가독성 향상",
                "qsort, signal 등 표준 라이브러리 활용"
            ])
        ],
        "code": r'''#include <stdio.h>
int add(int a, int b) { return a+b; }
int sub(int a, int b) { return a-b; }
int mul(int a, int b) { return a*b; }
void calculate(int x, int y,
    int (*op)(int,int), const char *name) {
    printf("%s(%d, %d) = %d\n",
           name, x, y, op(x,y));
}
int main() {
    int (*fp)(int, int);
    fp = add;
    printf("fp(3,4) = %d\n", fp(3,4));
    typedef int (*Operation)(int, int);
    Operation ops[] = {add, sub, mul};
    const char *names[] = {"add","sub","mul"};
    for(int i=0;i<3;i++)
        calculate(10, 3, ops[i], names[i]);
    return 0;
}''',
        "memory": [
            ("0x00401130", "add()", "(함수코드)", "function", "func"),
            ("0x00401160", "sub()", "(함수코드)", "function", "func"),
            ("0x00401190", "mul()", "(함수코드)", "function", "func"),
            ("0x7FFE0028", "fp", "0x00401130", "int(*)(int,int)", "ptr"),
            ("0x7FFE0030", "ops[0]", "0x00401130", "Operation", "ptr"),
            ("0x7FFE0038", "ops[1]", "0x00401160", "Operation", "ptr"),
            ("0x7FFE0040", "ops[2]", "0x00401190", "Operation", "ptr"),
        ],
        "arrows": "fp → add() / ops[0]→add, ops[1]→sub, ops[2]→mul : 함수 포인터가 코드 영역을 가리킴",
        "keypoints": [
            "함수명 = 함수의 시작 주소",
            "함수 포인터로 간접 호출 가능",
            "typedef로 가독성 향상",
            "콜백 패턴과 전략 패턴 구현에 활용"
        ]
    },
    {
        "id": 20,
        "title": "포인터 함수 vs 함수 포인터 비교",
        "concepts": [
            ("포인터 함수 (Pointer-returning Function)", [
                "int* func() : 포인터를 반환하는 함수",
                "함수가 동적 할당 주소 등을 반환",
                "호출자가 반환된 포인터의 메모리 관리 책임",
            ]),
            ("함수 포인터 (Function Pointer)", [
                "int (*func)() : 함수를 가리키는 포인터 변수",
                "괄호가 핵심 차이: (*func)의 괄호",
                "콜백, 전략 패턴, 함수 배열에 활용",
                "typedef로 가독성 향상 권장"
            ])
        ],
        "code": r'''#include <stdio.h>
#include <stdlib.h>
// 1. 포인터 함수: 포인터를 반환하는 함수
int* createValue(int val) {
    int *p = (int*)malloc(sizeof(int));
    *p = val;
    return p;
}
// 2. 함수 포인터: 함수를 가리키는 포인터
int add(int a, int b) { return a+b; }
int sub(int a, int b) { return a-b; }
int main() {
    // 포인터 함수 사용
    int *result = createValue(42);
    printf("포인터 함수 결과: %d\n", *result);
    free(result);
    // 함수 포인터 사용
    int (*op)(int, int) = add;
    printf("함수 포인터(add): %d\n", op(10, 3));
    op = sub;
    printf("함수 포인터(sub): %d\n", op(10, 3));
    // 함수 포인터 배열
    typedef int (*CalcFunc)(int, int);
    CalcFunc funcs[] = {add, sub};
    for(int i=0;i<2;i++)
        printf("funcs[%d](10,3) = %d\n",
               i, funcs[i](10,3));
    return 0;
}''',
        "memory": [
            ("0x00401130", "add()", "(함수코드)", "function", "func"),
            ("0x00401160", "sub()", "(함수코드)", "function", "func"),
            ("0x00401100", "createValue()", "(함수코드)", "function", "func"),
            ("0x00B71A40", "*result (힙)", "42", "int (동적)", "heap"),
            ("0x7FFE0038", "result (스택)", "0x00B71A40", "int*", "ptr"),
            ("0x7FFE0040", "op (스택)", "→ add/sub", "int(*)(int,int)", "ptr"),
        ],
        "arrows": "포인터 함수: createValue()→힙 주소 반환 / 함수 포인터: op→add/sub 코드 가리킴",
        "keypoints": [
            "int* func() = 포인터를 반환하는 함수",
            "int (*func)() = 함수를 가리키는 포인터",
            "괄호 유무가 핵심 차이",
            "포인터 함수: 동적 할당 결과 반환 / 함수 포인터: 콜백, 전략 패턴"
        ]
    },
]


# ============================================================
# HELPER FUNCTIONS
# ============================================================

def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_with_fill(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_name=FONT_KOREAN,
                font_size=14, font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                word_wrap=True):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_name=FONT_KOREAN, font_size=14,
                  font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                  level=0, space_before=Pt(4), space_after=Pt(4)):
    """Add a new paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    p.level = level
    p.space_before = space_before
    p.space_after = space_after
    return p


def create_section_title_slide(prs, lecture_id, title):
    """Create a section title slide with blue background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BLUE_PRIMARY)

    # Decorative top line
    add_shape_with_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), BLUE_ACCENT)

    # Lecture number - large
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.333), Inches(1.2),
                f"LECTURE {lecture_id:02d}",
                font_name=FONT_CODE, font_size=48, font_color=BLUE_ACCENT,
                bold=True, alignment=PP_ALIGN.CENTER)

    # Title
    add_textbox(slide, Inches(0.5), Inches(3.0), Inches(12.333), Inches(1.5),
                title,
                font_name=FONT_KOREAN, font_size=40, font_color=WHITE,
                bold=True, alignment=PP_ALIGN.CENTER)

    # Bottom decorative line
    add_shape_with_fill(slide, Inches(4), Inches(4.8), Inches(5.333), Inches(0.04), BLUE_ACCENT)

    # Subtitle
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12.333), Inches(0.6),
                "C언어 포인터 완전정복 - 64비트 시스템 기반",
                font_name=FONT_KOREAN, font_size=14, font_color=GRAY_TEXT,
                alignment=PP_ALIGN.CENTER)

    return slide


def create_concept_slide(prs, lecture_id, title, concepts):
    """Create concept explanation slide(s) with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Header bar
    add_shape_with_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), BLUE_PRIMARY)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6),
                f"Lecture {lecture_id:02d} | {title} - 핵심 개념",
                font_name=FONT_KOREAN, font_size=22, font_color=WHITE, bold=True)

    # Blue accent line
    add_shape_with_fill(slide, Inches(0), Inches(0.9), Inches(13.333), Inches(0.04), BLUE_ACCENT)

    # Content area
    y_pos = 1.2
    for section_title, bullets in concepts:
        # Section heading
        txBox = add_textbox(slide, Inches(0.8), Inches(y_pos), Inches(11.5), Inches(0.5),
                            f"\u25A0 {section_title}",
                            font_name=FONT_KOREAN, font_size=20, font_color=BLUE_PRIMARY,
                            bold=True)
        y_pos += 0.5

        # Bullet points
        for bullet in bullets:
            txBox2 = add_textbox(slide, Inches(1.3), Inches(y_pos), Inches(11), Inches(0.35),
                                 f"\u2022 {bullet}",
                                 font_name=FONT_KOREAN, font_size=15, font_color=GRAY_DARK)
            y_pos += 0.38

        y_pos += 0.15  # gap between sections

    return slide


def create_code_slide(prs, lecture_id, title, code):
    """Create a code example slide with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Header
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
                f"Lecture {lecture_id:02d} | {title} - 코드 예제",
                font_name=FONT_KOREAN, font_size=18, font_color=BLUE_ACCENT, bold=True)

    # Separator line
    add_shape_with_fill(slide, Inches(0.5), Inches(0.7), Inches(12.333), Inches(0.02), BLUE_ACCENT)

    # Code block background
    code_bg = add_shape_with_fill(slide, Inches(0.4), Inches(0.85),
                                  Inches(12.5), Inches(6.3),
                                  RGBColor(0x0F, 0x17, 0x2A))
    code_bg.line.color.rgb = RGBColor(0x33, 0x44, 0x55)
    code_bg.line.width = Pt(1)

    # Code text
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.1), Inches(6.1))
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = code.split('\n')
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line
        p.font.name = FONT_CODE
        p.font.color.rgb = WHITE
        p.space_before = Pt(0)
        p.space_after = Pt(0)

        # Calculate font size based on code length
        max_line_len = max(len(l) for l in lines)
        num_lines = len(lines)
        if num_lines > 28 or max_line_len > 60:
            p.font.size = Pt(10)
        elif num_lines > 22 or max_line_len > 50:
            p.font.size = Pt(11)
        elif num_lines > 16:
            p.font.size = Pt(12)
        else:
            p.font.size = Pt(13)

        # Simple syntax coloring via keywords
        stripped = line.strip()
        if stripped.startswith('//') or stripped.startswith('/*'):
            p.font.color.rgb = RGBColor(0x6A, 0x99, 0x55)  # green comments
        elif stripped.startswith('#include') or stripped.startswith('#define'):
            p.font.color.rgb = ORANGE_CODE
        elif any(stripped.startswith(kw) for kw in ['int ', 'char ', 'void ', 'float ', 'double ',
                                                     'long ', 'short ', 'const ', 'static ',
                                                     'typedef ', 'struct ']):
            p.font.color.rgb = CYAN_CODE
        elif 'printf' in stripped or 'malloc' in stripped or 'free' in stripped or 'calloc' in stripped or 'realloc' in stripped:
            p.font.color.rgb = YELLOW_CODE
        elif 'return' in stripped:
            p.font.color.rgb = ORANGE_CODE

    return slide


def get_row_color(mem_type):
    """Get background color based on memory type."""
    type_map = {
        "int": MEM_INT_BG,
        "ptr": MEM_PTR_BG,
        "char": MEM_CHAR_BG,
        "arr": MEM_ARR_BG,
        "heap": MEM_HEAP_BG,
        "func": MEM_FUNC_BG,
    }
    return type_map.get(mem_type, MEM_DEFAULT_BG)


def create_memory_slide(prs, lecture_id, title, memory_data, arrows):
    """Create a memory map slide with a colored table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Header bar
    add_shape_with_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), BLUE_PRIMARY)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6),
                f"Lecture {lecture_id:02d} | {title} - 메모리 맵",
                font_name=FONT_KOREAN, font_size=22, font_color=WHITE, bold=True)

    add_shape_with_fill(slide, Inches(0), Inches(0.9), Inches(13.333), Inches(0.04), BLUE_ACCENT)

    # Table
    rows_count = len(memory_data) + 1  # +1 for header
    cols_count = 4
    table_left = Inches(0.8)
    table_top = Inches(1.2)
    table_width = Inches(11.5)

    # Calculate row height based on number of rows
    max_table_height = 4.8
    row_height = min(0.5, max_table_height / rows_count)

    table_height = Inches(row_height * rows_count)

    table_shape = slide.shapes.add_table(rows_count, cols_count, table_left, table_top,
                                         table_width, table_height)
    table = table_shape.table

    # Set column widths
    col_widths = [Inches(2.2), Inches(3.0), Inches(3.3), Inches(3.0)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    headers = ["주소 (Address)", "변수명 (Name)", "값 (Value)", "타입 (Type)"]
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = MEM_HEADER_BG
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_KOREAN
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for row_idx, (addr, name, value, dtype, color_type) in enumerate(memory_data):
        row_color = get_row_color(color_type)
        data = [addr, name, value, dtype]
        for col_idx, cell_text in enumerate(data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color
            p = cell.text_frame.paragraphs[0]
            if col_idx == 0:  # Address column - monospace
                p.font.name = FONT_CODE
            elif col_idx == 2:  # Value column - monospace
                p.font.name = FONT_CODE
            else:
                p.font.name = FONT_KOREAN
            p.font.size = Pt(12)
            p.font.color.rgb = GRAY_DARK
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Arrow descriptions below table
    arrow_y = 1.2 + row_height * rows_count + 0.3

    # Legend box
    legend_bg = add_shape_with_fill(slide, Inches(0.8), Inches(arrow_y),
                                    Inches(11.5), Inches(0.45),
                                    RGBColor(0xF0, 0xF4, 0xFF))
    legend_bg.line.color.rgb = BLUE_LIGHT
    legend_bg.line.width = Pt(1)

    add_textbox(slide, Inches(1.0), Inches(arrow_y + 0.05), Inches(11.0), Inches(0.35),
                f"\u2192 포인터 관계: {arrows}",
                font_name=FONT_KOREAN, font_size=12, font_color=BLUE_PRIMARY, bold=False)

    # Color legend
    legend_y = arrow_y + 0.6
    legend_items = [
        (MEM_INT_BG, "int/기본형"),
        (MEM_PTR_BG, "포인터"),
        (MEM_CHAR_BG, "char"),
        (MEM_ARR_BG, "배열"),
        (MEM_HEAP_BG, "힙(동적)"),
        (MEM_FUNC_BG, "함수/코드"),
    ]

    x_offset = 0.8
    for color, label in legend_items:
        # Color box
        box = add_shape_with_fill(slide, Inches(x_offset), Inches(legend_y),
                                  Inches(0.25), Inches(0.25), color)
        box.line.color.rgb = GRAY_TEXT
        box.line.width = Pt(0.5)
        # Label
        add_textbox(slide, Inches(x_offset + 0.3), Inches(legend_y - 0.02),
                    Inches(1.3), Inches(0.3),
                    label, font_name=FONT_KOREAN, font_size=10, font_color=GRAY_DARK)
        x_offset += 1.8

    return slide


def create_keypoints_slide(prs, lecture_id, title, keypoints):
    """Create a key points summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Header bar
    add_shape_with_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), BLUE_PRIMARY)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6),
                f"Lecture {lecture_id:02d} | {title} - 핵심 정리",
                font_name=FONT_KOREAN, font_size=22, font_color=WHITE, bold=True)

    add_shape_with_fill(slide, Inches(0), Inches(0.9), Inches(13.333), Inches(0.04), BLUE_ACCENT)

    # Key points with checkmarks
    y_pos = 1.5
    for idx, point in enumerate(keypoints):
        # Checkmark circle
        circle = add_shape_with_fill(slide, Inches(1.2), Inches(y_pos + 0.05),
                                     Inches(0.4), Inches(0.4), BLUE_PRIMARY,
                                     MSO_SHAPE.OVAL)
        # Checkmark text inside circle
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = "\u2713"
        p.font.name = FONT_CODE
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Point text
        add_textbox(slide, Inches(1.8), Inches(y_pos), Inches(10.5), Inches(0.5),
                    point,
                    font_name=FONT_KOREAN, font_size=18, font_color=GRAY_DARK)

        # Subtle separator line
        if idx < len(keypoints) - 1:
            add_shape_with_fill(slide, Inches(1.8), Inches(y_pos + 0.55),
                                Inches(10), Inches(0.01), GRAY_LIGHT)

        y_pos += 0.7

    # Bottom accent
    add_shape_with_fill(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5), BLUE_PALE)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.333), Inches(0.4),
                f"\u2605 Lecture {lecture_id:02d} 핵심 포인트를 반드시 기억하세요!",
                font_name=FONT_KOREAN, font_size=13, font_color=BLUE_PRIMARY,
                bold=True, alignment=PP_ALIGN.CENTER)

    return slide


# ============================================================
# MAIN PRESENTATION GENERATION
# ============================================================

def generate_presentation():
    """Generate the complete PowerPoint presentation."""
    prs = Presentation()

    # Set slide size to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print("=" * 60)
    print("C언어 포인터 완전정복 - PPT 생성 시작")
    print("=" * 60)

    # ========================================
    # SLIDE 1: Opening Title Slide
    # ========================================
    print("[1/23] 표지 슬라이드 생성 중...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)

    # Top decorative bar
    add_shape_with_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_ACCENT)

    # Main title
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12.333), Inches(1.2),
                "C언어 포인터 완전정복",
                font_name=FONT_KOREAN, font_size=52, font_color=WHITE,
                bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Inches(0.5), Inches(2.9), Inches(12.333), Inches(0.8),
                "64비트 시스템 기반 메모리 시각화 학습",
                font_name=FONT_KOREAN, font_size=24, font_color=BLUE_ACCENT,
                alignment=PP_ALIGN.CENTER)

    # Decorative line
    add_shape_with_fill(slide, Inches(4), Inches(4.0), Inches(5.333), Inches(0.04), BLUE_LIGHT)

    # Description bullets
    desc_items = [
        "20개 강의 | 단계별 포인터 학습 커리큘럼",
        "코드 예제 + 메모리 맵으로 완벽 이해",
        "기초부터 함수 포인터까지 완전 정복"
    ]
    y = 4.5
    for item in desc_items:
        add_textbox(slide, Inches(2), Inches(y), Inches(9.333), Inches(0.45),
                    f"\u25B8 {item}",
                    font_name=FONT_KOREAN, font_size=16, font_color=GRAY_TEXT,
                    alignment=PP_ALIGN.CENTER)
        y += 0.45

    # Bottom bar
    add_shape_with_fill(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5), BLUE_PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.333), Inches(0.4),
                "StudyPointer | C Language Pointer Mastery Course",
                font_name=FONT_CODE, font_size=14, font_color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # ========================================
    # SLIDE 2: Table of Contents
    # ========================================
    print("[2/23] 목차 슬라이드 생성 중...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Header
    add_shape_with_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), DARK_NAVY)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6),
                "\u2630 학습 목차 - 전체 20강 커리큘럼",
                font_name=FONT_KOREAN, font_size=26, font_color=WHITE, bold=True)

    add_shape_with_fill(slide, Inches(0), Inches(0.9), Inches(13.333), Inches(0.04), BLUE_ACCENT)

    # Two-column layout for TOC
    col1_x = 0.6
    col2_x = 6.8
    y_start = 1.2
    row_height = 0.34

    for i, lec in enumerate(LECTURES):
        if i < 10:
            x = col1_x
            y = y_start + i * row_height
        else:
            x = col2_x
            y = y_start + (i - 10) * row_height

        # Number badge
        badge_color = BLUE_PRIMARY if i < 10 else RGBColor(0x15, 0x80, 0x3D)
        badge = add_shape_with_fill(slide, Inches(x), Inches(y),
                                    Inches(0.45), Inches(0.28), badge_color,
                                    MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = badge.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = f"{lec['id']:02d}"
        p.font.name = FONT_CODE
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Title text
        add_textbox(slide, Inches(x + 0.55), Inches(y - 0.02), Inches(5.5), Inches(0.32),
                    lec['title'],
                    font_name=FONT_KOREAN, font_size=13, font_color=GRAY_DARK)

    # Section labels
    add_textbox(slide, Inches(col1_x), Inches(y_start + 10 * row_height + 0.1),
                Inches(5.5), Inches(0.3),
                "\u25B2 기초 과정 (Lecture 01-10)",
                font_name=FONT_KOREAN, font_size=11, font_color=BLUE_PRIMARY, bold=True)

    add_textbox(slide, Inches(col2_x), Inches(y_start + 10 * row_height + 0.1),
                Inches(5.5), Inches(0.3),
                "\u25B2 심화 과정 (Lecture 11-20)",
                font_name=FONT_KOREAN, font_size=11, font_color=RGBColor(0x15, 0x80, 0x3D), bold=True)

    # Bottom note
    add_shape_with_fill(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5), BLUE_PALE)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.333), Inches(0.4),
                "각 강의는 개념 설명 + 코드 예제 + 메모리 맵 + 핵심 정리로 구성되어 있습니다",
                font_name=FONT_KOREAN, font_size=13, font_color=BLUE_PRIMARY,
                alignment=PP_ALIGN.CENTER)

    # ========================================
    # LECTURES 1-20: Generate slides for each
    # ========================================
    for lec in LECTURES:
        lid = lec['id']
        ltitle = lec['title']
        print(f"[{lid + 2}/23] Lecture {lid:02d}: {ltitle} 생성 중...")

        # 1. Section title slide
        create_section_title_slide(prs, lid, ltitle)

        # 2. Concept slide(s)
        create_concept_slide(prs, lid, ltitle, lec['concepts'])

        # 3. Code example slide
        create_code_slide(prs, lid, ltitle, lec['code'])

        # 4. Memory map slide
        create_memory_slide(prs, lid, ltitle, lec['memory'], lec['arrows'])

        # 5. Key points slide
        create_keypoints_slide(prs, lid, ltitle, lec['keypoints'])

    # ========================================
    # CLOSING SLIDE
    # ========================================
    print("[23/23] 마무리 슬라이드 생성 중...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)

    # Top bar
    add_shape_with_fill(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_ACCENT)

    # Completion message
    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12.333), Inches(1.0),
                "\u2705",
                font_name=FONT_CODE, font_size=60, font_color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.0),
                "학습 완료!",
                font_name=FONT_KOREAN, font_size=48, font_color=WHITE,
                bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.5), Inches(3.5), Inches(12.333), Inches(0.8),
                "C++ 포인터로의 여정을 시작하세요",
                font_name=FONT_KOREAN, font_size=28, font_color=BLUE_ACCENT,
                alignment=PP_ALIGN.CENTER)

    # Decorative line
    add_shape_with_fill(slide, Inches(4), Inches(4.6), Inches(5.333), Inches(0.04), BLUE_LIGHT)

    # Summary stats
    stats = [
        "20개 강의 완료 | 100+ 슬라이드",
        "포인터 기초부터 함수 포인터까지",
        "64비트 메모리 시각화 학습 완성"
    ]
    y = 5.0
    for stat in stats:
        add_textbox(slide, Inches(2), Inches(y), Inches(9.333), Inches(0.4),
                    f"\u2606 {stat}",
                    font_name=FONT_KOREAN, font_size=16, font_color=GRAY_TEXT,
                    alignment=PP_ALIGN.CENTER)
        y += 0.4

    # Bottom bar
    add_shape_with_fill(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5), BLUE_PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.333), Inches(0.4),
                "Thank you! | StudyPointer - C Pointer Mastery Complete",
                font_name=FONT_CODE, font_size=14, font_color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # ========================================
    # SAVE
    # ========================================
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "C언어_포인터_완전정복.pptx")
    prs.save(output_path)

    total_slides = len(prs.slides)
    print()
    print("=" * 60)
    print(f"생성 완료!")
    print(f"  파일: {output_path}")
    print(f"  총 슬라이드 수: {total_slides}")
    print(f"  강의 수: 20개")
    print(f"  슬라이드 구성: 표지 + 목차 + (20 x 5 강의 슬라이드) + 마무리")
    print("=" * 60)


if __name__ == "__main__":
    generate_presentation()
